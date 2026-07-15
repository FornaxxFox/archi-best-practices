"""
pptx_helpers.py · 小红书图文工厂 v1.3.3 共享工具

集中：
1. 调色板常量 (PALETTE_MAP) — 海报生成 + 封面插图共享
2. OVERLAY_COLOR — 封面图叠加色，唯一 source of truth
3. embed_cover_image() — 封面图嵌入函数
4. _normalize_previous_episodes() — 兼容 v1.3.0/v1.3.1/v1.3.2/v1.3.3 四种 previous_episodes 形状

M29 修复：海报生成 Step 3b + 封面插图 之前各有一份重复的 python-pptx 代码，常量还冲突。
v1.3.3 起所有 pptx 相关代码必须 import 此模块。
"""

from __future__ import annotations

import re
from typing import Any

try:
    from pptx.util import Inches, Emu
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
except ImportError:
    raise ImportError("需要 python-pptx: pip install --user python-pptx")


# === 调色板常量 (v1.3.3 唯一 source of truth) ===

PALETTE_MAP: dict[str, dict[str, RGBColor]] = {
    "暖调知识风": {
        "bg":     RGBColor(0xFA, 0xF5, 0xEB),
        "fg":     RGBColor(0x3A, 0x2A, 0x1A),
        "accent": RGBColor(0xC8, 0x6B, 0x2E),
    },
    "清新种草风": {
        "bg":     RGBColor(0xFF, 0xFD, 0xF7),
        "fg":     RGBColor(0x2D, 0x4A, 0x2D),
        "accent": RGBColor(0xF8, 0xB4, 0xC8),
    },
    "深色hook风": {
        "bg":     RGBColor(0x14, 0x1B, 0x2E),
        "fg":     RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0xFF, 0x6B, 0x35),
    },
    "文艺情绪风": {
        "bg":     RGBColor(0xF5, 0xE8, 0xC8),
        "fg":     RGBColor(0x4A, 0x33, 0x1F),
        "accent": RGBColor(0x7A, 0x8B, 0x4F),
    },
}

# 唯一 source of truth（v1.3.3 之前 海报生成 = 0x1A1A2E, 封面插图 = 0x000000，冲突）
OVERLAY_COLOR: RGBColor = RGBColor(0x1A, 0x1A, 0x2E)


def resolve_palette(style_hint: str) -> dict[str, RGBColor]:
    """从 style_hint 字符串解析为 palette dict。fallback = 暖调知识风。"""
    return PALETTE_MAP.get(style_hint, PALETTE_MAP["暖调知识风"])


# === 兼容性 normalizer (v1.3.3 · HIGH-4 修复) ===

_VOL_RE = re.compile(r"[-/]vol(\d+)$")


def normalize_previous_episodes(previous_episodes: list | None) -> list[int]:
    """兼容 v1.3.0/v1.3.1 (字符串路径) / v1.3.2 (对象数组) / 空 三种形状，返回 done_vols: list[int]。

    Args:
        previous_episodes: 原始 previous_episodes 列表，可能是 None / 字符串列表 / 对象列表

    Returns:
        list[int]: 已发布期数（不含当前期）

    Examples:
        >>> normalize_previous_episodes(["projects/拆书-vol1", "projects/拆书-vol2"])
        [1, 2]
        >>> normalize_previous_episodes([{"project_id": "拆书", "episode": 3}])
        [3]
        >>> normalize_previous_episodes(None)
        []
        >>> normalize_previous_episodes(["projects/no-vol-suffix"])
        []
    """
    if not previous_episodes:
        return []
    done_vols: list[int] = []
    for ep in previous_episodes:
        if isinstance(ep, str):
            m = _VOL_RE.search(ep)
            if m:
                done_vols.append(int(m.group(1)))
            # 无 -volN 后缀 → 跳过（无法确定 vol 编号）
        elif isinstance(ep, dict):
            try:
                done_vols.append(int(ep.get("episode", 0)))
            except (TypeError, ValueError):
                continue
    return done_vols


# === 封面图嵌入函数 (v1.3.3 M29 修复) ===

def embed_cover_image(
    slide,
    image_path: str,
    placement: str = "background",
    title: str = "",
    subtitle: str = "",
    palette: dict[str, RGBColor] | None = None,
) -> None:
    """在 slide 上嵌入封面图 + 文字叠加（v1.3.3 共享版本）。

    Args:
        slide: python-pptx slide object
        image_path: 封面图绝对/相对路径
        placement: "background" | "top_half" | "center_float"
        title: 主标题
        subtitle: 副标题
        palette: PALETTE_MAP 中的 palette dict；缺省 = 暖调知识风
    """
    import os
    palette = palette or PALETTE_MAP["暖调知识风"]

    if not image_path or not os.path.exists(image_path):
        # 无封面图：纯色底
        from pptx.oxml.ns import qn
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = palette["bg"]
        return

    if placement == "background":
        slide.shapes.add_picture(
            image_path,
            left=Emu(0), top=Emu(0),
            width=Inches(7.5), height=Inches(10.0),
        )
        # 底部暗色遮罩确保标题可读 — 共享 OVERLAY_COLOR
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Emu(0), top=Inches(6.5),
            width=Inches(7.5), height=Inches(3.5),
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = OVERLAY_COLOR
        overlay.line.fill.background()

    elif placement == "top_half":
        slide.shapes.add_picture(
            image_path,
            left=Emu(0), top=Emu(0),
            width=Inches(7.5), height=Inches(5.5),
        )
    # else (center_float): 留空，由调用者自行添加


# === A/B 变体维护 (v1.3.3 M19) ===

def format_ab_output(variant: str, output_path: str) -> dict[str, str]:
    """生成 ab_outputs 数组项的规范格式。"""
    from datetime import datetime, timezone, timedelta
    tz = timezone(timedelta(hours=8))
    return {
        "variant": variant,
        "output_path": output_path,
        "generated_at": datetime.now(tz).isoformat(timespec="seconds"),
    }


# === v1.3.3 自检（import 时跑） ===

if __name__ == "__main__":
    # Smoke test
    assert normalize_previous_episodes(None) == []
    assert normalize_previous_episodes([]) == []
    assert normalize_previous_episodes(["projects/拆书-深度工作-vol1", "projects/拆书-深度工作-vol2"]) == [1, 2]
    assert normalize_previous_episodes(["projects/拆书-深度工作/vol3"]) == [3]  # /vol also works
    assert normalize_previous_episodes([{"project_id": "拆书", "episode": 5}, {"episode": 3}]) == [5, 3]
    assert normalize_previous_episodes([{"episode": "bad"}]) == []
    assert normalize_previous_episodes(["no-suffix"]) == []
    assert resolve_palette("暖调知识风") is not None
    assert resolve_palette("unknown") is not None  # fallback
    assert OVERLAY_COLOR is not None
    print("pptx_helpers.py self-test OK")
